from dotenv import load_dotenv
import streamlit as st
from PyPDF2 import PdfReader
from langchain.text_splitter import CharacterTextSplitter
from langchain.embeddings.openai import OpenAIEmbeddings
from langchain.vectorstores import FAISS
from langchain.chains.question_answering import load_qa_chain
from langchain.llms import OpenAI
from langchain.callbacks import get_openai_callback
import os
import json
from render import user_msg_container_html_template, bot_msg_container_html_template
from langchain.chains import ConversationalRetrievalChain
from langchain.chat_models import ChatOpenAI
from langchain.vectorstores import Chroma
import matplotlib.pyplot as plt
import seaborn as sns
import re
import time
import plotly.graph_objects as go
import plotly.io as pio
import tempfile
import glob
import pandas as pd
from langchain.agents import create_csv_agent
from langchain.agents.agent_types import AgentType
from pandasai import PandasAI
import docx
from pptx import Presentation

def submit():
    st.session_state.input = st.session_state.widget
    st.session_state.widget = ''

def get_pdf_text(pdf_docs):
    text = ""
    for pdf in pdf_docs:
        pdf_reader = PdfReader(pdf)
        for page in pdf_reader.pages:
            text += page.extract_text()
    return text
    
    
def get_docx_text(word_docs):
    text = ""
    for word in word_docs:
        doc_reader = docx.Document(word)
        for paragraph in doc_reader.paragraphs:
            text += paragraph.text
    return text
    
def get_ppt_text(ppt_docs):
    text = ""
    for ppt in ppt_docs:
        ppt_reader = Presentation(ppt)
        for slide in ppt_reader.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text
    return text
    
    
def get_text_chunks(text):
    text_splitter = CharacterTextSplitter(
        separator="\n",
        chunk_size=1000,
        chunk_overlap=200,
        length_function=len
    )
    chunks = text_splitter.split_text(text)
    return chunks
    
def get_vectorstore(text_chunks):
    embeddings = OpenAIEmbeddings()
    # embeddings = HuggingFaceInstructEmbeddings(model_name="hkunlp/instructor-xl")
    vectorstore = FAISS.from_texts(texts=text_chunks, embedding=embeddings)
    return vectorstore

    
def main():
    load_dotenv()
    current_directory = os.path.dirname(os.path.abspath(__file__))
    css_file = os.path.join(current_directory, "css.txt")
    st.set_page_config(page_title="Ask your documents",initial_sidebar_state="expanded")
    st.set_option('deprecation.showPyplotGlobalUse', False)
    # --- LOAD CSS ---
    with open(css_file) as f:
        st.markdown(f"<style>{f.read()}</style>", unsafe_allow_html=True)
    st.markdown("<h1 class='center-header'>Ask your documents and .CSVs💬</h1>", unsafe_allow_html=True)
    # upload file
    uploaded_files = st.file_uploader(
            "Upload your PDFs, Word docs, PowerPoints or CSVs here", accept_multiple_files=True)
    # Create a temporary directory to store the uploaded files
    temp_dir = tempfile.mkdtemp()
    temp_path = temp_dir

    # Save the uploaded files to the temporary directory
    for i, uploaded_file in enumerate(uploaded_files):
        # Get the filename from the uploaded file
        filename = uploaded_file.name

        # Extract the file extension using os.path.splitext
        file_extension = os.path.splitext(filename)[1]
        file_path = os.path.join(temp_path, filename)
        with open(file_path, "wb") as f:
            f.write(uploaded_file.getvalue())

    # Get all CSV files in the temporary directory
    csv_files = glob.glob(os.path.join(temp_path, "*.csv"))
    pdf_docs = glob.glob(os.path.join(temp_path, "*.pdf"))
    word_docs = glob.glob(os.path.join(temp_path, "*.docx"))
    ppt_docs = glob.glob(os.path.join(temp_path, "*.ppt*"))

    # Create an empty dictionary to store the DataFrames
    dataframes = {}

    # Iterate over the CSV files
    for file in csv_files:
        # Extract the file name without extension
        file_name = file.split('/')[-1].split('.')[0]
        
        # Read the CSV file into a DataFrame
        df = pd.read_csv(file)
        
        # Store the DataFrame in the dictionary
        dataframes[file_name] = df
            
    pdf_raw_text = get_pdf_text(pdf_docs)
    docx_raw_text = get_docx_text(word_docs)
    ppt_raw_text = get_ppt_text(ppt_docs)
    
    raw_text = pdf_raw_text+docx_raw_text+ppt_raw_text

    # get the text chunks
    text_chunks = get_text_chunks(raw_text)
    
    if text_chunks:
    # create vector store
        knowledge_base = get_vectorstore(text_chunks)

            
    # Initialize session state for chat history
    if 'chat_history' not in st.session_state:
        st.session_state.chat_history = []
    if 'chat_history_docs' not in st.session_state:
        st.session_state.chat_history_docs = []
    if 'input' not in st.session_state:
        st.session_state.input = ''
    
    c1 = st.columns(2)
    
    AI_MODEL_OPTIONS = [
    "gpt-3.5-turbo",
    "gpt-4",
    "gpt-4-32k"
]
    selected_option = st.selectbox("Select an LLM model", AI_MODEL_OPTIONS)
    
    if  selected_option:
        user_question = st.text_input("Ask your question:", key='widget', on_change=submit)

        if st.session_state.input:

            if selected_option == "gpt-4":
                if raw_text and not csv_files:
                    docs = knowledge_base.similarity_search(st.session_state.input)
                
                    qa = ConversationalRetrievalChain.from_llm(
                        ChatOpenAI(temperature=0.1, model="gpt-4"),
                        knowledge_base.as_retriever()
                    )
                    with get_openai_callback() as cb:
                      response = qa({"question": st.session_state.input, "chat_history": st.session_state.chat_history_docs})
                      st.session_state.chat_history_docs.append((st.session_state.input, response["answer"]))
                      st.session_state.chat_history.append((st.session_state.input, response["answer"]))
                      
                elif csv_files and not raw_text:
                    llm = OpenAI()
                    # create PandasAI object, passing the LLM
                    pandas_ai = PandasAI(llm, conversational=False, verbose=True)
                    #pandas_ai.clear_cache()
                
                    if any(word in st.session_state.input for word in ["plot","chart","Plot","Chart"]):
                        question = st.session_state.input + ' ' + 'using seaborn'
                    else:
                        question = st.session_state.input

                    #fig = go.Figure()
                    fig = plt.gcf()
                    x = pandas_ai.run(list(dataframes.values()), question)

                    if fig.get_axes():
                        st.session_state.chat_history.append((st.session_state.input, fig))
                    
                    else:
                        st.session_state.chat_history.append((st.session_state.input, x))
                    
                else:
                    try:
                        docs = knowledge_base.similarity_search(st.session_state.input)
                    
                        qa = ConversationalRetrievalChain.from_llm(
                            ChatOpenAI(temperature=0.1, model="gpt-4"),
                            knowledge_base.as_retriever()
                        )
                        with get_openai_callback() as cb:
                          response = qa({"question": st.session_state.input, "chat_history": st.session_state.chat_history_docs})
                          st.session_state.chat_history_docs.append((st.session_state.input, response["answer"]))
                          st.session_state.chat_history.append((st.session_state.input, response["answer"]))
                          
                    except:
                        llm = OpenAI()
                        # create PandasAI object, passing the LLM
                        pandas_ai = PandasAI(llm, conversational=False, verbose=True)
                        #pandas_ai.clear_cache()
                    
                        if any(word in st.session_state.input for word in ["plot","chart","Plot","Chart"]):
                            question = st.session_state.input + ' ' + 'using seaborn'
                        else:
                            question = st.session_state.input

                        fig = plt.gcf()
                        x = pandas_ai.run(list(dataframes.values()), question)

                        if fig.get_axes():
                            st.session_state.chat_history.append((st.session_state.input, fig))
                        
                        else:
                            st.session_state.chat_history.append((st.session_state.input, x))
                            st.session_state.chat_history_docs.append((st.session_state.input, str(x)))
                        
            elif selected_option == "gpt-3.5-turbo":
                if raw_text and not csv_files:
                    docs = knowledge_base.similarity_search(st.session_state.input)
                
                    qa = ConversationalRetrievalChain.from_llm(
                        ChatOpenAI(temperature=0.1, model="gpt-3.5-turbo"),
                        knowledge_base.as_retriever()
                    )
                    with get_openai_callback() as cb:
                      response = qa({"question": st.session_state.input, "chat_history": st.session_state.chat_history_docs})
                      st.session_state.chat_history_docs.append((st.session_state.input, response["answer"]))
                      st.session_state.chat_history.append((st.session_state.input, response["answer"]))
                      
                elif csv_files and not raw_text:
                    llm = OpenAI()
                    # create PandasAI object, passing the LLM
                    pandas_ai = PandasAI(llm, conversational=False, verbose=True)
                    #pandas_ai.clear_cache()
                
                    if any(word in st.session_state.input for word in ["plot","chart","Plot","Chart"]):
                        question = st.session_state.input + ' ' + 'using seaborn'
                    else:
                        question = st.session_state.input

                    fig = plt.gcf()
                    x = pandas_ai.run(list(dataframes.values()), question)

                    if fig.get_axes():
                        st.session_state.chat_history.append((st.session_state.input, fig))
                    
                    else:
                        st.session_state.chat_history.append((st.session_state.input, x))
                else:
                    try:
                        docs = knowledge_base.similarity_search(st.session_state.input)
                    
                        qa = ConversationalRetrievalChain.from_llm(
                            ChatOpenAI(temperature=0.1, model="gpt-3.5-turbo"),
                            knowledge_base.as_retriever()
                        )
                        with get_openai_callback() as cb:
                          response = qa({"question": st.session_state.input, "chat_history": st.session_state.chat_history_docs})
                          st.session_state.chat_history_docs.append((st.session_state.input, response["answer"]))
                          st.session_state.chat_history.append((st.session_state.input, response["answer"]))
                          
                    except:
                        llm = OpenAI()
                        # create PandasAI object, passing the LLM
                        pandas_ai = PandasAI(llm, conversational=False, verbose=True)
                        #pandas_ai.clear_cache()
                    
                        if any(word in st.session_state.input for word in ["plot","chart","Plot","Chart"]):
                            question = st.session_state.input + ' ' + 'using seaborn'
                        else:
                            question = st.session_state.input

                        fig = plt.gcf()
                        x = pandas_ai.run(list(dataframes.values()), question)

                        if fig.get_axes():
                            st.session_state.chat_history.append((st.session_state.input, fig))
                        
                        else:
                            st.session_state.chat_history.append((st.session_state.input, x))
                            st.session_state.chat_history_docs.append((st.session_state.input, str(x)))
                            
            else:
                if raw_text and not csv_files:
                    docs = knowledge_base.similarity_search(st.session_state.input)
                
                    qa = ConversationalRetrievalChain.from_llm(
                        ChatOpenAI(temperature=0.1, model="gpt-4"),
                        knowledge_base.as_retriever()
                    )
                    with get_openai_callback() as cb:
                      response = qa({"question": st.session_state.input, "chat_history": st.session_state.chat_history_docs})
                      st.session_state.chat_history_docs.append((st.session_state.input, response["answer"]))
                      st.session_state.chat_history.append((st.session_state.input, response["answer"]))
                      
                elif csv_files and not raw_text:
                    llm = OpenAI()
                    # create PandasAI object, passing the LLM
                    pandas_ai = PandasAI(llm, conversational=False, verbose=True)
                    #pandas_ai.clear_cache()
                
                    if any(word in st.session_state.input for word in ["plot","chart","Plot","Chart"]):
                        question = st.session_state.input + ' ' + 'using seaborn'
                    else:
                        question = st.session_state.input

                    #fig = go.Figure()
                    fig = plt.gcf()
                    x = pandas_ai.run(list(dataframes.values()), question)

                    if fig.get_axes():
                        st.session_state.chat_history.append((st.session_state.input, fig))
                    
                    else:
                        st.session_state.chat_history.append((st.session_state.input, x))
            if selected_option == "gpt-4-32k":
                if raw_text and not csv_files:
                    docs = knowledge_base.similarity_search(st.session_state.input)
                
                    qa = ConversationalRetrievalChain.from_llm(
                        ChatOpenAI(temperature=0.1, model="gpt-4"),
                        knowledge_base.as_retriever()
                    )
                    with get_openai_callback() as cb:
                      response = qa({"question": st.session_state.input, "chat_history": st.session_state.chat_history_docs})
                      st.session_state.chat_history_docs.append((st.session_state.input, response["answer"]))
                      st.session_state.chat_history.append((st.session_state.input, response["answer"]))
                      
                elif csv_files and not raw_text:
                    llm = OpenAI()
                    # create PandasAI object, passing the LLM
                    pandas_ai = PandasAI(llm, conversational=False, verbose=True)
                    #pandas_ai.clear_cache()
                
                    if any(word in st.session_state.input for word in ["plot","chart","Plot","Chart"]):
                        question = st.session_state.input + ' ' + 'using seaborn'
                    else:
                        question = st.session_state.input

                    #fig = go.Figure()
                    fig = plt.gcf()
                    x = pandas_ai.run(list(dataframes.values()), question)

                    if fig.get_axes():
                        st.session_state.chat_history.append((st.session_state.input, fig))
                    
                    else:
                        st.session_state.chat_history.append((st.session_state.input, x))
                    
                else:
                    try:
                        docs = knowledge_base.similarity_search(st.session_state.input)
                    
                        qa = ConversationalRetrievalChain.from_llm(
                            ChatOpenAI(temperature=0.1, model="gpt-4"),
                            knowledge_base.as_retriever()
                        )
                        with get_openai_callback() as cb:
                          response = qa({"question": st.session_state.input, "chat_history": st.session_state.chat_history_docs})
                          st.session_state.chat_history_docs.append((st.session_state.input, response["answer"]))
                          st.session_state.chat_history.append((st.session_state.input, response["answer"]))
                          
                    except:
                        llm = OpenAI()
                        # create PandasAI object, passing the LLM
                        pandas_ai = PandasAI(llm, conversational=False, verbose=True)
                        #pandas_ai.clear_cache()
                    
                        if any(word in st.session_state.input for word in ["plot","chart","Plot","Chart"]):
                            question = st.session_state.input + ' ' + 'using seaborn'
                        else:
                            question = st.session_state.input

                        fig = plt.gcf()
                        x = pandas_ai.run(list(dataframes.values()), question)

                        if fig.get_axes():
                            st.session_state.chat_history.append((st.session_state.input, fig))
                        
                        else:
                            st.session_state.chat_history.append((st.session_state.input, x))
                            st.session_state.chat_history_docs.append((st.session_state.input, str(x)))

            st.session_state.input = ''
 


    # Display chat history
    for message in st.session_state.chat_history[::-1]:
        if message[0]:
            st.write(f"<div class='custom-text'>{user_msg_container_html_template.replace('$MSG', message[0])}</div>", unsafe_allow_html=True)
            try:
                st.pyplot(message[1])
            except:
                st.write(f"<div class='custom-text'>{bot_msg_container_html_template.replace('$MSG', str(message[1]))}</div>", unsafe_allow_html=True)
            
if __name__ == '__main__':
    main()

